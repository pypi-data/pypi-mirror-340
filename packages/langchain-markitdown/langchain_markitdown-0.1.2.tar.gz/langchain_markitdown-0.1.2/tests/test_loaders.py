import pytest
from markitdown_langchain import (
    DocxLoader,
    PptxLoader,
    XlsxLoader,
    BaseMarkitdownLoader  # Import BaseMarkitdownLoader
)
from langchain_core.documents import Document
import os
from unittest.mock import patch

# --- Fixtures ---
@pytest.fixture(scope="module")
def test_docx_file(tmpdir_factory):
    """Creates a temporary DOCX file for testing."""
    fn = tmpdir_factory.mktemp("data").join("test.docx")
    # Create a minimal DOCX file (you might need to install python-docx)
    try:
        from docx import Document as DocxDocument
        doc = DocxDocument()
        doc.add_paragraph("This is a test document.")
        doc.save(fn)
    except ImportError:
        pytest.skip("docx package not installed. Install with 'pip install python-docx'")
    return str(fn)

@pytest.fixture(scope="module")
def test_pptx_file(tmpdir_factory):
    """Creates a temporary PPTX file for testing."""
    fn = tmpdir_factory.mktemp("data").join("test.pptx")
    # Create a minimal PPTX file (you might need to install python-pptx)
    try:
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "Test Presentation"
        prs.save(fn)
    except ImportError:
        pytest.skip("pptx package not installed. Install with 'pip install python-pptx'")
    return str(fn)

@pytest.fixture(scope="module")
def test_xlsx_file(tmpdir_factory):
    """Creates a temporary XLSX file for testing."""
    fn = tmpdir_factory.mktemp("data").join("test.xlsx")
    # Create a minimal XLSX file (you might need to install openpyxl)
    try:
        import openpyxl
        wb = openpyxl.Workbook()
        ws = wb.active
        ws['A1'] = "Test Data"
        wb.save(fn)
    except ImportError:
        pytest.skip("openpyxl package not installed. Install with 'pip install openpyxl'")
    return str(fn)

# --- Test Cases ---

def test_docx_loader(test_docx_file):
    """Test loading a DOCX file."""
    loader = DocxLoader(test_docx_file)
    documents = loader.load()

    assert isinstance(documents, list)
    assert len(documents) > 0
    assert isinstance(documents[0], Document)
    assert "test.docx" in documents[0].metadata["source"]

def test_pptx_loader(test_pptx_file):
    """Test loading a PPTX file."""
    loader = PptxLoader(test_pptx_file)
    documents = loader.load()

    assert isinstance(documents, list)
    assert len(documents) > 0
    assert isinstance(documents[0], Document)
    assert "test.pptx" in documents[0].metadata["source"]

def test_xlsx_loader(test_xlsx_file):
    """Test loading a XLSX file."""
    loader = XlsxLoader(test_xlsx_file)
    documents = loader.load()

    assert isinstance(documents, list)
    assert len(documents) > 0
    assert isinstance(documents[0], Document)
    assert "test.xlsx" in documents[0].metadata["source"]

def test_base_loader_file_not_found():
    """Test handling of non-existent file in BaseMarkitdownLoader."""
    with pytest.raises(ValueError) as excinfo:  # Changed from FileNotFoundError to ValueError
        BaseMarkitdownLoader("non_existent_file.txt").load()
    assert "Markitdown conversion failed" in str(excinfo.value)
    assert "File not found" in str(excinfo.value)

def test_base_loader_invalid_file():
    """Test handling of an invalid file type with BaseMarkitdownLoader."""
    with patch("os.path.getsize") as mock_getsize, pytest.raises(ValueError) as excinfo:
        mock_getsize.side_effect = FileNotFoundError
        loader = BaseMarkitdownLoader("invalid_file.xyz")
        loader.load()
    assert "Markitdown conversion failed" in str(excinfo.value)