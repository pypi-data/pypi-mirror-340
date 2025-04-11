from typing import List, Dict, Any, Optional
from langchain_core.documents import Document
from .base_loader import BaseMarkitdownLoader
from langchain_text_splitters import MarkdownHeaderTextSplitter
from langchain_core.language_models import BaseChatModel
import re
import os
from .utils import get_image_caption  # Import the function

class PptxLoader(BaseMarkitdownLoader):
    def __init__(self, file_path: str, split_by_page: bool = False, llm: Optional[BaseChatModel] = None):
        super().__init__(file_path)
        self.split_by_page = split_by_page
        self.llm = llm

    def load(
        self, 
        headers_to_split_on: Optional[List[str]] = None
    ) -> List[Document]:
        """Load a PPTX file and convert it to Langchain documents, splitting by Markdown headers."""
        try:
            # Basic converter for fallback when MarkitdownConverterOptions isn't available
            from markitdown import MarkItDown
            
            # Create basic metadata
            metadata: Dict[str, Any] = {
                "source": self.file_path,
                "file_name": self._get_file_name(self.file_path),
                "file_size": self._get_file_size(self.file_path),
                "conversion_success": True,
            }
            
            # Extract detailed metadata from PPTX using python-pptx
            try:
                from pptx import Presentation
                prs = Presentation(self.file_path)
                
                # Basic presentation stats
                metadata["slide_count"] = len(prs.slides)
                
                # Core properties
                if hasattr(prs, 'core_properties'):
                    core_props = prs.core_properties
                    if hasattr(core_props, 'author') and core_props.author:
                        metadata["author"] = core_props.author
                    if hasattr(core_props, 'title') and core_props.title:
                        metadata["title"] = core_props.title
                    if hasattr(core_props, 'subject') and core_props.subject:
                        metadata["subject"] = core_props.subject
                    if hasattr(core_props, 'keywords') and core_props.keywords:
                        metadata["keywords"] = core_props.keywords
                    if hasattr(core_props, 'created') and core_props.created:
                        metadata["created"] = str(core_props.created)
                    if hasattr(core_props, 'modified') and core_props.modified:
                        metadata["modified"] = str(core_props.modified)
                    if hasattr(core_props, 'last_modified_by') and core_props.last_modified_by:
                        metadata["last_modified_by"] = core_props.last_modified_by
                    if hasattr(core_props, 'category') and core_props.category:
                        metadata["category"] = core_props.category
                    if hasattr(core_props, 'revision') and core_props.revision:
                        metadata["revision"] = core_props.revision
                
                # Count media elements
                image_count = 0
                text_box_count = 0
                chart_count = 0
                table_count = 0
                
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'shape_type'):
                            if shape.shape_type == 13:  # Picture
                                image_count += 1
                            elif shape.shape_type == 17:  # TextBox
                                text_box_count += 1
                            elif shape.shape_type == 3:  # Chart
                                chart_count += 1
                            elif shape.shape_type == 19:  # Table
                                table_count += 1
                
                metadata["image_count"] = image_count
                metadata["text_box_count"] = text_box_count
                metadata["chart_count"] = chart_count
                metadata["table_count"] = table_count
                
            except Exception as e:
                # If metadata extraction fails, continue with basic metadata
                metadata["metadata_extraction_error"] = str(e)
            
            # Try to use MarkitdownConverterOptions for image captioning if available
            try:
                from markitdown import MarkitdownConverterOptions
                converter = MarkItDown(options=MarkitdownConverterOptions(
                    llm_for_image_caption=lambda file_stream, stream_info, **kwargs: 
                        get_image_caption(self.llm, file_stream, stream_info) if self.llm else None
                ))
            except ImportError:
                # Fall back to basic converter if MarkitdownConverterOptions is not available
                converter = MarkItDown()
                metadata["image_captioning"] = "disabled - MarkitdownConverterOptions not available"
            
            # Convert the presentation to markdown
            result = converter.convert(self.file_path)

            # Define default headers to split on if not provided
            if headers_to_split_on is None:
                headers_to_split_on = [
                    ("#", "Header 1"),
                    ("##", "Header 2"),
                    ("###", "Header 3"),
                ]

            if self.split_by_page:
                # Split by slide number indicators
                documents = []
                slide_pattern = r"\n\n<!-- Slide number: (\d+) -->\n"
                slide_splits = re.split(slide_pattern, result.text_content)
                
                # The first element will be the content before the first slide indicator
                current_page_content = slide_splits[0]
                current_page_num = 1  # Assume the content before the first indicator belongs to slide 1

                for i in range(1, len(slide_splits), 2):
                    if current_page_content.strip():  # Avoid empty pages
                        page_metadata = metadata.copy()
                        page_metadata["page_number"] = current_page_num
                        page_metadata["content_type"] = "presentation_slide"
                        
                        # Split page content by headers
                        markdown_splitter = MarkdownHeaderTextSplitter(headers_to_split_on=headers_to_split_on)
                        page_splits = markdown_splitter.split_text(current_page_content)
                        
                        # Add split documents with updated metadata
                        for split in page_splits:
                            split.metadata.update(page_metadata)
                            documents.append(split)

                    current_page_num = int(slide_splits[i])  # Get the slide number from the indicator
                    current_page_content = slide_splits[i + 1]  # Get the content of the current slide

                # Add the last page
                if current_page_content.strip():
                    page_metadata = metadata.copy()
                    page_metadata["page_number"] = current_page_num
                    page_metadata["content_type"] = "presentation_slide"
                    markdown_splitter = MarkdownHeaderTextSplitter(headers_to_split_on=headers_to_split_on)
                    page_splits = markdown_splitter.split_text(current_page_content)
                    for split in page_splits:
                        split.metadata.update(page_metadata)
                        documents.append(split)
            else:
                # If not splitting by page, perform header-based splitting on the entire document
                markdown_splitter = MarkdownHeaderTextSplitter(headers_to_split_on=headers_to_split_on)
                documents = markdown_splitter.split_text(result.text_content)
                for doc in documents:
                    doc.metadata.update(metadata)
                    doc.metadata["content_type"] = "presentation_full"

            return documents

        except Exception as e:
            raise ValueError(f"Failed to load and convert PPTX file: {e}")
