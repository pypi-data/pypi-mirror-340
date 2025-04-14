import magic
from pathlib import Path
from typing import Optional, Dict, Any, Union, List
import fitz  # PyMuPDF
import docx
from PIL import Image
from file_prompt_cli.gcp.gemini_client import GeminiClient
from file_prompt_cli.core.preprocessor import DocumentPreprocessor
import os
import logging
import base64
import mimetypes
import pandas as pd
import json
from pptx import Presentation

logger = logging.getLogger(__name__)

class FileProcessor:
    """Handles file processing and content extraction."""
    
    SUPPORTED_MIME_TYPES = {
        'application/pdf': 'pdf',
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document': 'docx',
        'text/plain': 'txt',
        'image/jpeg': 'jpg',
        'image/png': 'png',
        'image/webp': 'webp',
        'text/csv': 'csv',
        'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': 'xlsx',
        'application/vnd.ms-excel': 'xls',
        'application/vnd.ms-excel.sheet.macroEnabled.12': 'xlsm',
        'application/vnd.openxmlformats-officedocument.spreadsheetml.template': 'xltx',
        'application/vnd.ms-excel.template.macroEnabled.12': 'xltm',
        'application/msword': 'doc',
        'application/vnd.openxmlformats-officedocument.presentationml.presentation': 'pptx',
        'application/vnd.ms-powerpoint': 'ppt',
        'application/vnd.openxmlformats-officedocument.presentationml.slideshow': 'ppsx',
        'application/vnd.ms-powerpoint.slideshow.macroEnabled.12': 'ppsm',
        'application/vnd.openxmlformats-officedocument.presentationml.template': 'potx',
        'application/vnd.ms-powerpoint.template.macroEnabled.12': 'potm'
    }
    
    def __init__(self, gemini_client: Optional[GeminiClient] = None):
        """Initialize the FileProcessor.
        
        Args:
            gemini_client: Optional GeminiClient instance. If not provided,
                          a new one will be created.
        """
        self.gemini_client = gemini_client or GeminiClient()
        self._preprocessors = {
            "pdf": self._preprocess_pdf,
            "docx": self._preprocess_docx,
            "txt": self._preprocess_txt,
            "csv": self._preprocess_csv,
            "ppt": self._preprocess_ppt,
            "pptx": self._preprocess_ppt,
            "image": self._preprocess_image
        }
        self.preprocessor = DocumentPreprocessor()
        logger.info("FileProcessor initialized with DocumentPreprocessor and GeminiClient")
    
    def detect_file_type(self, file_path: Path) -> str:
        """Detect the MIME type of a file.
        
        Args:
            file_path: Path to the file
            
        Returns:
            The detected MIME type
        """
        try:
            mime = magic.Magic(mime=True)
            mime_type = mime.from_file(str(file_path))
            
            if mime_type not in self.SUPPORTED_MIME_TYPES:
                raise ValueError(f"Unsupported file type: {mime_type}")
                
            return mime_type
            
        except Exception as e:
            logger.error(f"Error detecting file type: {str(e)}")
            raise
    
    def _detect_mime_type(self, file_path: str) -> str:
        """Detect MIME type of file."""
        try:
            mime_type, _ = mimetypes.guess_type(file_path)
            if not mime_type:
                # Try to determine from extension
                ext = os.path.splitext(file_path)[1].lower()
                if ext in ['.pdf']:
                    mime_type = 'application/pdf'
                elif ext in ['.docx']:
                    mime_type = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
                elif ext in ['.txt']:
                    mime_type = 'text/plain'
                elif ext in ['.jpg', '.jpeg']:
                    mime_type = 'image/jpeg'
                elif ext in ['.png']:
                    mime_type = 'image/png'
                elif ext in ['.webp']:
                    mime_type = 'image/webp'
                elif ext in ['.csv']:
                    mime_type = 'text/csv'
                elif ext in ['.xlsx']:
                    mime_type = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
                elif ext in ['.xls']:
                    mime_type = 'application/vnd.ms-excel'
                elif ext in ['.ppt']:
                    mime_type = 'application/vnd.ms-powerpoint'
                elif ext in ['.pptx']:
                    mime_type = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
            
            logger.info(f"Detected MIME type for {file_path}: {mime_type}")
            return mime_type
        except Exception as e:
            logger.error(f"Error detecting MIME type: {str(e)}")
            return 'application/octet-stream'

    def _preprocess_csv(self, file_path: str) -> Dict[str, Any]:
        """Preprocess CSV file."""
        try:
            logger.info(f"Preprocessing CSV file: {file_path}")
            
            # Read CSV file
            df = pd.read_csv(file_path)
            logger.info(f"Successfully read CSV file. Shape: {df.shape}")
            logger.debug(f"Columns: {df.columns.tolist()}")
            logger.debug(f"First few rows:\n{df.head()}")
            
            # Get basic statistics
            stats = {
                "rows": len(df),
                "columns": len(df.columns),
                "column_names": list(df.columns),
                "data_types": df.dtypes.astype(str).to_dict(),
                "sample_data": df.head().to_dict('records')
            }
            logger.debug(f"Generated stats: {json.dumps(stats, indent=2)}")
            
            # Convert to markdown table for analysis
            markdown_table = df.head().to_markdown()
            logger.debug(f"Generated markdown table:\n{markdown_table}")
            
            result = {
                "status": "success",
                "content": markdown_table,
                "metadata": {
                    "file_type": "csv",
                    "stats": stats,
                    "is_ready_for_gemini": True
                }
            }
            
            logger.info(f"CSV preprocessing complete: {result['status']}")
            return result
            
        except Exception as e:
            logger.error(f"Error preprocessing CSV: {str(e)}")
            logger.exception("Full traceback:")
            return {
                "status": "error",
                "error": str(e),
                "metadata": {
                    "file_type": "csv",
                    "is_ready_for_gemini": False
                }
            }

    def _preprocess_excel(self, file_path: str) -> Dict[str, Any]:
        """Preprocess Excel file."""
        try:
            logger.info(f"Preprocessing Excel file: {file_path}")
            
            # Read Excel file
            excel_file = pd.ExcelFile(file_path)
            sheets = excel_file.sheet_names
            
            # Process each sheet
            sheet_data = {}
            for sheet in sheets:
                df = pd.read_excel(excel_file, sheet_name=sheet)
                
                # Get basic statistics for each sheet
                stats = {
                    "rows": len(df),
                    "columns": len(df.columns),
                    "column_names": list(df.columns),
                    "data_types": df.dtypes.astype(str).to_dict(),
                    "sample_data": df.head().to_dict('records')
                }
                
                # Convert to markdown table for analysis
                markdown_table = df.head().to_markdown()
                
                sheet_data[sheet] = {
                    "stats": stats,
                    "markdown": markdown_table
                }
            
            result = {
                "status": "success",
                "content": "\n\n".join([f"## Sheet: {sheet}\n\n{data['markdown']}" 
                                   for sheet, data in sheet_data.items()]),
                "metadata": {
                    "file_type": "excel",
                    "sheets": sheets,
                    "sheet_data": sheet_data,
                    "is_ready_for_gemini": True
                }
            }
            
            logger.info(f"Excel preprocessing complete: {result['status']}")
            return result
            
        except Exception as e:
            logger.error(f"Error preprocessing Excel: {str(e)}")
            return {
                "status": "error",
                "error": str(e),
                "metadata": {
                    "file_type": "excel",
                    "is_ready_for_gemini": False
                }
            }

    def _preprocess_pdf(self, file_path: str) -> Dict[str, Any]:
        """Preprocess PDF file."""
        try:
            logger.info(f"Preprocessing PDF file: {file_path}")
            
            # Open PDF file
            doc = fitz.open(file_path)
            
            # Extract text from all pages
            text = ""
            for page in doc:
                text += page.get_text()
            
            # Get metadata
            metadata = doc.metadata
            
            result = {
                "status": "success",
                "content": text,
                "metadata": {
                    "file_type": "pdf",
                    "page_count": len(doc),
                    "metadata": metadata,
                    "is_ready_for_gemini": True
                }
            }
            
            logger.info(f"PDF preprocessing complete: {result['status']}")
            return result
            
        except Exception as e:
            logger.error(f"Error preprocessing PDF: {str(e)}")
            logger.exception("Full traceback:")
            return {
                "status": "error",
                "error": str(e),
                "metadata": {
                    "file_type": "pdf",
                    "is_ready_for_gemini": False
                }
            }

    def _preprocess_txt(self, file_path: str) -> Dict[str, Any]:
        """Preprocess text file."""
        try:
            logger.info(f"Preprocessing text file: {file_path}")
            
            # Read text file
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
            
            result = {
                "status": "success",
                "content": text,
                "metadata": {
                    "file_type": "txt",
                    "is_ready_for_gemini": True
                }
            }
            
            logger.info(f"Text preprocessing complete: {result['status']}")
            return result
            
        except Exception as e:
            logger.error(f"Error preprocessing text file: {str(e)}")
            logger.exception("Full traceback:")
            return {
                "status": "error",
                "error": str(e),
                "metadata": {
                    "file_type": "txt",
                    "is_ready_for_gemini": False
                }
            }

    def _preprocess_docx(self, file_path: str) -> Dict[str, Any]:
        """Preprocess DOCX file."""
        try:
            logger.info(f"Preprocessing DOCX file: {file_path}")
            
            # Open DOCX file
            doc = docx.Document(file_path)
            
            # Extract text from all paragraphs
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            
            result = {
                "status": "success",
                "content": text,
                "metadata": {
                    "file_type": "docx",
                    "is_ready_for_gemini": True
                }
            }
            
            logger.info(f"DOCX preprocessing complete: {result['status']}")
            return result
            
        except Exception as e:
            logger.error(f"Error preprocessing DOCX: {str(e)}")
            logger.exception("Full traceback:")
            return {
                "status": "error",
                "error": str(e),
                "metadata": {
                    "file_type": "docx",
                    "is_ready_for_gemini": False
                }
            }

    def _preprocess_image(self, file_path: str) -> Dict[str, Any]:
        """Preprocess image file."""
        try:
            logger.info(f"Preprocessing image file: {file_path}")
            
            # Open image file
            with Image.open(file_path) as img:
                # Get image metadata
                metadata = {
                    "format": img.format,
                    "mode": img.mode,
                    "size": img.size,
                    "info": img.info
                }
                
                # Read image file as bytes and encode to base64
                with open(file_path, 'rb') as f:
                    image_bytes = f.read()
                    image_base64 = base64.b64encode(image_bytes).decode('utf-8')
            
            result = {
                "status": "success",
                "content": image_base64,  # Pass the base64 string directly
                "metadata": {
                    "file_type": "image",
                    "image_metadata": metadata,
                    "is_ready_for_gemini": True
                }
            }
            
            logger.info(f"Image preprocessing complete: {result['status']}")
            return result
            
        except Exception as e:
            logger.error(f"Error preprocessing image: {str(e)}")
            logger.exception("Full traceback:")
            return {
                "status": "error",
                "error": str(e),
                "metadata": {
                    "file_type": "image",
                    "is_ready_for_gemini": False
                }
            }

    def _preprocess_ppt(self, file_path: str) -> Dict[str, Any]:
        """Preprocess PowerPoint file.
        
        Args:
            file_path: Path to the PowerPoint file
            
        Returns:
            Dict containing preprocessing results
        """
        try:
            logger.info(f"Preprocessing PowerPoint file: {file_path}")
            
            # Open the presentation
            prs = Presentation(file_path)
            
            # Extract text from all slides
            text_content = []
            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                if slide_text:
                    text_content.append("\n".join(slide_text))
            
            # Get presentation metadata
            metadata = {
                "file_type": "ppt",
                "slide_count": len(prs.slides),
                "is_ready_for_gemini": True
            }
            
            # Combine all text content
            content = "\n\n".join(text_content)
            
            result = {
                "status": "success",
                "content": content,
                "metadata": metadata
            }
            
            logger.info(f"PowerPoint preprocessing complete: {result['status']}")
            return result
            
        except Exception as e:
            logger.error(f"Error preprocessing PowerPoint: {str(e)}")
            logger.exception("Full traceback:")
            return {
                "status": "error",
                "error": str(e),
                "metadata": {
                    "file_type": "ppt",
                    "is_ready_for_gemini": False
                }
            }

    def preprocess_file(self, file_path: str, prompt: Optional[str] = None, metadata: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
        """Preprocess a file and analyze its content.
        
        Args:
            file_path: Path to the file to process
            prompt: Optional custom prompt for content analysis
            metadata: Optional metadata to pass to the analysis
            
        Returns:
            Dict containing preprocessing results and analysis
        """
        try:
            logger.info(f"Starting preprocessing for file: {file_path}")
            
            # Detect file type
            mime_type = self._detect_mime_type(file_path)
            logger.info(f"Detected MIME type: {mime_type}")
            
            # Preprocess based on file type
            if mime_type == 'application/pdf':
                result = self._preprocess_pdf(file_path)
            elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
                result = self._preprocess_docx(file_path)
            elif mime_type == 'text/plain':
                result = self._preprocess_txt(file_path)
            elif mime_type in ['image/jpeg', 'image/png', 'image/webp']:
                result = self._preprocess_image(file_path)
            elif mime_type == 'text/csv':
                result = self._preprocess_csv(file_path)
            elif mime_type in ['application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 
                              'application/vnd.ms-excel',
                              'application/vnd.ms-excel.sheet.macroEnabled.12',
                              'application/vnd.openxmlformats-officedocument.spreadsheetml.template',
                              'application/vnd.ms-excel.template.macroEnabled.12']:
                result = self._preprocess_excel(file_path)
            elif mime_type in ['application/msword',
                              'application/vnd.openxmlformats-officedocument.presentationml.presentation',
                              'application/vnd.ms-powerpoint',
                              'application/vnd.openxmlformats-officedocument.presentationml.slideshow',
                              'application/vnd.ms-powerpoint.slideshow.macroEnabled.12',
                              'application/vnd.openxmlformats-officedocument.presentationml.template',
                              'application/vnd.ms-powerpoint.template.macroEnabled.12']:
                result = self._preprocess_ppt(file_path)
            else:
                raise ValueError(f"Unsupported file type: {mime_type}")
            
            # Add file metadata
            result['file_type'] = mime_type
            result['file_path'] = file_path
            
            # Analyze content with Gemini
            if result['status'] == 'success' and result['metadata']['is_ready_for_gemini']:
                analysis_prompt = prompt or "Analyze the following content and provide insights:"
                analysis = self.gemini_client.analyze_content(
                    content=result['content'],
                    prompt=analysis_prompt,
                    metadata=result['metadata']
                )
                
                # Use the Gemini analysis response directly as the content
                result['content'] = analysis.get('response', 'No analysis available')
            
            # Remove metadata before returning
            if 'metadata' in result:
                del result['metadata']
            
            logger.info(f"Preprocessing completed for file: {file_path}")
            return result
            
        except Exception as e:
            logger.error(f"Error preprocessing file {file_path}: {str(e)}")
            logger.exception("Full traceback:")
            return {
                "status": "error",
                "error": str(e),
                "file_type": mime_type
            }
